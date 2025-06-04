from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

# Filter image files and sort
def get_sorted_images(folder: str) -> list:
  return sorted([
    os.path.join(folder, f)
    for f in os.listdir(folder)
    if f.lower().endswith(('.png', '.jpg', '.jpeg', '.bmp', '.gif'))
  ])

# Normalize: removes prefix like "1-" and extension
def normalize(name) -> str:
    name = os.path.splitext(name)[0]
    if '-' in name and name.split('-')[0].isdigit():
        name = name.split('-', 1)[1]
    return name.strip().lower()

def create_powerpoint() -> Presentation:
  prs = Presentation()
  prs.slide_width = Inches(12.01)
  prs.slide_height = Inches(8.47)
  return prs

def create_slide(prs: Presentation) -> any:
  # Add a new slide (Title and Content layout)
  slide_layout = prs.slide_layouts[6]  # Use blank layout
  slide = prs.slides.add_slide(slide_layout)

  # Add background
  background = slide.background
  fill = background.fill
  fill.solid()
  fill.fore_color.rgb = RGBColor(0, 0, 0)  # Black
  return slide

def add_title(slide: any, text: str):
  title_shape = slide.shapes.add_textbox(Inches(0), Inches(7.18), Inches(12.01), Inches(0.59))
  text_frame = title_shape.text_frame
  p = text_frame.paragraphs[0]
  p.alignment = PP_ALIGN.CENTER
  run = p.add_run()
  run.text = text
  font = run.font
  font.name = 'Comic Sans MS'
  font.size = Pt(40)
  font.color.rgb = RGBColor(0xB9, 0xB4, 0x53)  # Hex #B9B453

def main():
  picture_top = Inches(0.69)
  picture_height = Inches(6.4)

  for group in [501, 502, 503, 504, 505, 506, 507, 508, 551, 552, 553, 554]:
    folder_left = f"../{group} Photos Initiales"
    folder_right = f"../{group}"

    left_images = get_sorted_images(folder_left)
    right_images = get_sorted_images(folder_right)
    
    left_map = {normalize(os.path.basename(f)): f for f in left_images}
    right_map = {normalize(os.path.basename(f)): f for f in right_images}

    # Keys to generate slides for: all in right (even if left is missing)
    all_keys = sorted(right_map.keys())
    
    prs = create_powerpoint()

    for key in all_keys:
      left_path = left_map.get(key)
      right_path = right_map.get(key)

      slide = create_slide(prs)
      slide_title = os.path.splitext(os.path.basename(right_path))[0]
      add_title(slide, slide_title)

      if left_path and right_path:
        # Add left image
        slide.shapes.add_picture(left_path, Inches(0.8), picture_top, height=picture_height)

        # Add right image
        slide.shapes.add_picture(right_path, Inches(6.4), picture_top, height=picture_height)

      elif right_path:
        # Add right image
        slide.shapes.add_picture(right_path, Inches(3.605), picture_top, height=picture_height)

    # Save the presentation
    prs.save(f"{group}.pptx")

if __name__ == "__main__":
  main()